"""
Flat-function PPTX builder. No classes, no wrapper stacking.
Driven by JSON config, supports image placement with captions.

Usage:
  python pptx_builder.py --config layout.json --output deck.pptx
  python pptx_builder.py --stdin < layout.json
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- Palettes ----
PALETTES = {
    "premium": {
        "text": "#161A1D", "muted": "#596168", "bg": "#F6F3EC",
        "panel": "#FFFCF5", "primary": "#1A3A5C", "accent": "#D66A2E",
        "accent2": "#245C4D", "white": "#FFFFFF", "border": "#E6DED1",
        "font": "Microsoft YaHei",
    },
    "modern": {
        "text": "#1A1A2E", "muted": "#6B7280", "bg": "#FFFFFF",
        "panel": "#F8F9FA", "primary": "#2563EB", "accent": "#7C3AED",
        "accent2": "#0EA5E9", "white": "#FFFFFF", "border": "#E5E7EB",
        "font": "Microsoft YaHei",
    },
    "dark": {
        "text": "#F0F0F0", "muted": "#9CA3AF", "bg": "#111827",
        "panel": "#1F2937", "primary": "#1D4ED8", "accent": "#F59E0B",
        "accent2": "#10B981", "white": "#FFFFFF", "border": "#374151",
        "font": "Microsoft YaHei",
    },
}


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_palette(config: dict) -> dict:
    """Build resolved palette dict with RGBColor values."""
    name = config.get("palette", "premium")
    raw = dict(PALETTES.get(name, PALETTES["premium"]))
    raw.update(config.get("palette_overrides", {}))
    return {k: hex_to_rgb(v) if k != "font" else v for k, v in raw.items()}


# ---- Core slide helpers (flat functions, no self) ----

def init_prs(config: dict) -> tuple:
    """Create Presentation, set dimensions, return (prs, palette, font)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    C = resolve_palette(config)
    font = C.pop("font", "Microsoft YaHei")
    return prs, C, font


def new_slide(prs, C) -> object:
    """Create a blank slide with full-slide background rectangle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["bg"]
    bg.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = bg._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)
    return slide


def add_tb(slide, x, y, w, h, fill=None):
    """Add a text box with standard margins and word wrap."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.15)
    box.text_frame.margin_right = Inches(0.15)
    box.text_frame.margin_top = Inches(0.1)
    box.text_frame.margin_bottom = Inches(0.1)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    return box


def set_p(p, text, color, size, bold=False, align=PP_ALIGN.LEFT, space=None,
          font="Microsoft YaHei"):
    """Format a paragraph."""
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    if space:
        p.space_after = Pt(space)


def add_p(tf, text, color, size, bold=False, align=PP_ALIGN.LEFT, space=None,
          font="Microsoft YaHei"):
    """Add a new formatted paragraph to text frame."""
    p = tf.add_paragraph()
    set_p(p, text, color, size, bold, align, space, font)
    return p


def add_footer(slide, page_num, title, C, font):
    """Add page number footer."""
    t = add_tb(slide, 0.75, 7.05, 3, 0.25)
    set_p(t.text_frame.paragraphs[0],
          f"{title} / {page_num:02d}", C["muted"], 10, font=font)


def add_rect(slide, x, y, w, h, color):
    """Add a filled rectangle with no outline."""
    r = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r


def add_rrect(slide, x, y, w, h, fill_color, border_color=None):
    """Add a rounded rectangle with fill and optional border."""
    r = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    r.fill.solid()
    r.fill.fore_color.rgb = fill_color
    if border_color:
        r.line.color.rgb = border_color
        r.line.width = Pt(0.5)
    else:
        r.line.fill.background()
    return r


def add_oval(slide, x, y, w, h, fill_color):
    """Add a filled oval."""
    o = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    o.fill.solid()
    o.fill.fore_color.rgb = fill_color
    o.line.fill.background()
    return o


def place_image(prs, slide, img_cfg, C, font):
    """Place an image on a slide with optional caption.

    img_cfg keys: path (required), layout, caption, x, y, w, h
    Layouts: background, left-half, right-half, top-half, bottom-half, center, inline
    """
    img_path = img_cfg["path"]
    if not Path(img_path).exists():
        print(f"WARNING: Image not found: {img_path}", file=sys.stderr)
        return

    layout = img_cfg.get("layout", "right-half")
    caption = img_cfg.get("caption", "")
    sw = prs.slide_width
    sh = prs.slide_height

    # For now, read image and compute position
    from PIL import Image
    im = Image.open(img_path)
    iw, ih = im.size

    # Compute position based on layout
    m = int(Inches(0.4))  # margin in EMU
    hw = sw // 2

    if layout == "background":
        l, t, w, h = 0, 0, sw, sh
    elif layout == "left-half":
        region_w, region_h = hw - m * 2, sh - m * 2
        scale = min(region_w / (iw or 1), region_h / (ih or 1))
        w, h = int(iw * scale), int(ih * scale)
        l, t = m + (region_w - w) // 2, m + (region_h - h) // 2
    elif layout == "right-half":
        region_w, region_h = hw - m * 2, sh - m * 2
        scale = min(region_w / (iw or 1), region_h / (ih or 1))
        w, h = int(iw * scale), int(ih * scale)
        l = hw + m + (region_w - w) // 2
        t = m + (region_h - h) // 2
    elif layout == "top-half":
        region_w, region_h = sw - m * 2, sh // 2 - m * 2
        scale = min(region_w / (iw or 1), region_h / (ih or 1))
        w, h = int(iw * scale), int(ih * scale)
        l, t = m + (region_w - w) // 2, m + (region_h - h) // 2
    elif layout == "bottom-half":
        region_w, region_h = sw - m * 2, sh // 2 - m * 2
        scale = min(region_w / (iw or 1), region_h / (ih or 1))
        w, h = int(iw * scale), int(ih * scale)
        l = m + (region_w - w) // 2
        t = sh // 2 + m + (region_h - h) // 2
    elif layout == "center":
        region_w, region_h = sw - m * 2, sh - m * 2
        scale = min(region_w / (iw or 1), region_h / (ih or 1))
        w, h = int(iw * scale), int(ih * scale)
        l, t = m + (region_w - w) // 2, m + (region_h - h) // 2
    elif layout == "inline":
        l = int(Inches(img_cfg.get("x", 1)))
        t = int(Inches(img_cfg.get("y", 1)))
        w = int(Inches(img_cfg.get("w", 5)))
        h = int(Inches(img_cfg.get("h", 3)))
    else:
        region_w, region_h = hw - m * 2, sh - m * 2
        scale = min(region_w / (iw or 1), region_h / (ih or 1))
        w, h = int(iw * scale), int(ih * scale)
        l = hw + m + (region_w - w) // 2
        t = m + (region_h - h) // 2

    # Insert image
    pic = slide.shapes.add_picture(img_path, l, t, w, h)

    # Send behind text
    sp_tree = slide.shapes._spTree
    sp = pic._element
    sp_tree.remove(sp)
    # Insert after last text box (roughly behind text)
    children = list(sp_tree)
    sp_tree.insert(len(children) - 1, sp)

    # Add caption if provided
    if caption:
        cap_y = Inches(0.05)  # below image
        if layout in ("top-half",):
            cap_y = Inches(0.05)
        cap_box = add_tb(slide, l / 914400, (t + h) / 914400 + 0.08, w / 914400, 0.3)
        set_p(cap_box.text_frame.paragraphs[0], caption, C["muted"], 10,
              align=PP_ALIGN.CENTER, font=font)

    return pic


# ---- Slide builders (one per type) ----

def build_cover(prs, cfg, C, font, title, page_num):
    """Cover slide."""
    slide = new_slide(prs, C)
    add_rect(slide, 0, 0, 0.12, 7.5, C["primary"])

    eyebrow = cfg.get("eyebrow", "")
    if eyebrow:
        t = add_tb(slide, 0.75, 2.0, 11.8, 0.4)
        set_p(t.text_frame.paragraphs[0], eyebrow, C["primary"], 13, True, font=font)

    t = add_tb(slide, 0.75, 2.5, 11.8, 1.2)
    set_p(t.text_frame.paragraphs[0], cfg.get("title", title), C["text"], 42, True, font=font)

    subtitle = cfg.get("subtitle", "")
    if subtitle:
        t = add_tb(slide, 0.75, 3.8, 11.8, 0.5)
        set_p(t.text_frame.paragraphs[0], subtitle, C["muted"], 18, font=font)

    add_rect(slide, 0.75, 4.5, 1.2, 0.05, C["accent"])

    kicker = cfg.get("kicker", "")
    if kicker:
        t = add_tb(slide, 0.75, 5.2, 3, 0.3)
        set_p(t.text_frame.paragraphs[0], kicker, C["muted"], 12, font=font)

    add_footer(slide, page_num, title, C, font)


def build_section_header(prs, cfg, C, font, title, page_num):
    """Section divider slide."""
    slide = new_slide(prs, C)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.2), prs.slide_width, Inches(3.1)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = C["primary"]
    band.line.fill.background()

    t = add_tb(slide, 1.5, 2.7, 10.333, 1.0, C["primary"])
    set_p(t.text_frame.paragraphs[0], cfg.get("title", ""), C["white"], 40, True,
          PP_ALIGN.CENTER, font=font)

    subtitle = cfg.get("subtitle", "")
    if subtitle:
        t = add_tb(slide, 1.5, 3.8, 10.333, 0.5, C["primary"])
        set_p(t.text_frame.paragraphs[0], subtitle, C["white"], 18, False,
              PP_ALIGN.CENTER, font=font)

    add_footer(slide, page_num, title, C, font)


def build_content(prs, cfg, C, font, title, page_num):
    """Content slide: title + body + bullets + optional image."""
    slide = new_slide(prs, C)
    num_label = cfg.get("number", page_num)
    num_str = f"{num_label:02d}" if isinstance(num_label, int) else str(num_label)

    t = add_tb(slide, 0.75, 0.4, 1, 0.4)
    set_p(t.text_frame.paragraphs[0], num_str, C["primary"], 13, True, font=font)

    t = add_tb(slide, 0.75, 0.9, 11.8, 0.7)
    set_p(t.text_frame.paragraphs[0], cfg.get("title", ""), C["text"], 34, True, font=font)

    add_rect(slide, 0.75, 1.6, 1.2, 0.04, C["accent"])

    # Image on the right if specified
    has_image = "image" in cfg
    text_right_margin = 6.5 if has_image else 11.8  # narrower text area if image present

    y = 1.9
    body = cfg.get("body", "")
    if body:
        t = add_tb(slide, 0.75, y, text_right_margin, 0.5)
        set_p(t.text_frame.paragraphs[0], body, C["muted"], 15, font=font)
        y += 0.6

    bullets = cfg.get("bullets", [])
    if bullets:
        t = add_tb(slide, 0.75, y, text_right_margin, 5.5 - y)
        for i, b in enumerate(bullets):
            if i == 0:
                set_p(t.text_frame.paragraphs[0], f"  {b}", C["text"], 18, space=12, font=font)
            else:
                add_p(t.text_frame, f"  {b}", C["text"], 18, space=12, font=font)

    # Place image
    if has_image:
        img = cfg["image"]
        img.setdefault("layout", "right-half")
        place_image(prs, slide, img, C, font)

    add_footer(slide, page_num, title, C, font)


def build_cards(prs, cfg, C, font, title, page_num):
    """Cards slide: 2-3 horizontal cards."""
    slide = new_slide(prs, C)
    items = cfg.get("items", [])
    n = len(items)
    if n < 2:
        n = 2

    if cfg.get("title"):
        t = add_tb(slide, 0.75, 0.4, 11.8, 0.7)
        set_p(t.text_frame.paragraphs[0], cfg["title"], C["text"], 32, True, font=font)

    gap = 0.3
    cw = (11.8 - gap * (n - 1)) / n
    ch = 4.8
    cy = 1.3

    for i, item in enumerate(items):
        x = 0.75 + i * (cw + gap)
        add_rrect(slide, x, cy, cw, ch, C["panel"], C["border"])

        icon = item.get("icon", f"{i+1:02d}")
        o = add_oval(slide, x + 0.3, cy + 0.3, 0.6, 0.6, C["primary"])
        set_p(o.text_frame.paragraphs[0], icon, C["white"], 18, True, PP_ALIGN.CENTER, font=font)
        o.text_frame.margin_top = Inches(0.12)

        t = add_tb(slide, x + 0.3, cy + 1.1, cw - 0.6, 0.4)
        set_p(t.text_frame.paragraphs[0], item.get("title", ""), C["text"], 18, True, font=font)

        body = item.get("body", "")
        if body:
            t = add_tb(slide, x + 0.3, cy + 1.6, cw - 0.6, ch - 1.9)
            set_p(t.text_frame.paragraphs[0], body, C["muted"], 14, font=font)

        add_rect(slide, x + cw - 1.5, cy + ch - 0.06, 1.2, 0.04, C["accent"])

    add_footer(slide, page_num, title, C, font)


def build_side_by_side(prs, cfg, C, font, title, page_num):
    """Side-by-side: image one side, text the other."""
    slide = new_slide(prs, C)

    if cfg.get("title"):
        t = add_tb(slide, 0.75, 0.4, 11.8, 0.6)
        set_p(t.text_frame.paragraphs[0], cfg["title"], C["text"], 30, True, font=font)

    image_side = cfg.get("image_side", "right")
    image_path = cfg.get("image", {}).get("path", "") if isinstance(cfg.get("image"), dict) else cfg.get("image", "")

    img_w = 5.5
    img_h = 4.5
    img_y = 1.4
    img_x = 0.75 if image_side == "left" else (13.333 - 0.75 - img_w)
    text_x = 0.75 if image_side == "right" else (img_x + img_w + 0.5)
    # Ensure text doesn't overlap image
    if image_side == "right":
        text_w = img_x - text_x - 0.3
    else:
        text_w = 13.333 - text_x - 0.75

    if image_path and Path(image_path).exists():
        pic = slide.shapes.add_picture(
            str(image_path), Inches(img_x), Inches(img_y), Inches(img_w), Inches(img_h)
        )
        # Caption
        if isinstance(cfg.get("image"), dict) and cfg["image"].get("caption"):
            cap = add_tb(slide, img_x, img_y + img_h + 0.08, img_w, 0.3)
            set_p(cap.text_frame.paragraphs[0], cfg["image"]["caption"], C["muted"], 10,
                  align=PP_ALIGN.CENTER, font=font)
    else:
        ph = add_rrect(slide, img_x, img_y, img_w, img_h, C["panel"], C["border"])
        t = add_tb(slide, img_x + img_w / 2 - 1, img_y + img_h / 2 - 0.2, 2, 0.4)
        set_p(t.text_frame.paragraphs[0], "[ Image ]", C["muted"], 13, False, PP_ALIGN.CENTER, font=font)

    bullets = cfg.get("bullets", [])
    if bullets:
        t = add_tb(slide, text_x, img_y, text_w, img_h)
        for i, b in enumerate(bullets):
            if i == 0:
                set_p(t.text_frame.paragraphs[0], f"  {b}", C["text"], 18, space=14, font=font)
            else:
                add_p(t.text_frame, f"  {b}", C["text"], 18, space=14, font=font)

    add_footer(slide, page_num, title, C, font)


def build_comparison(prs, cfg, C, font, title, page_num):
    """Two-column comparison slide."""
    slide = new_slide(prs, C)

    if cfg.get("title"):
        t = add_tb(slide, 0.75, 0.4, 11.8, 0.6)
        set_p(t.text_frame.paragraphs[0], cfg["title"], C["text"], 30, True, font=font)

    labels = cfg.get("labels", ["Before", "After"])
    col_w, col_h, col_y = 5.5, 4.5, 1.3

    for side, data, label, is_primary in [
        ("left", cfg.get("left", {}), labels[0] if labels else "Before", False),
        ("right", cfg.get("right", {}), labels[1] if len(labels) > 1 else "After", True),
    ]:
        x = 0.75 if side == "left" else (13.333 - 0.75 - col_w)
        add_rrect(slide, x, col_y, col_w, col_h,
                  C["primary"] if is_primary else C["panel"])

        label_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.3), Inches(col_y + 0.3), Inches(1.2), Inches(0.35)
        )
        label_bg.fill.solid()
        label_bg.fill.fore_color.rgb = C["accent"] if is_primary else C["primary"]
        label_bg.line.fill.background()
        lc = C["text"] if is_primary else C["white"]
        set_p(label_bg.text_frame.paragraphs[0], label.upper(), lc, 11, True,
              PP_ALIGN.CENTER, font=font)
        label_bg.text_frame.margin_top = Inches(0.05)

        tc = C["white"] if is_primary else C["text"]
        ct = data.get("title", "")
        if ct:
            t = add_tb(slide, x + 0.3, col_y + 0.9, col_w - 0.6, 0.4)
            set_p(t.text_frame.paragraphs[0], ct, tc, 20, True, font=font)

        cb = data.get("body", "")
        if cb:
            t = add_tb(slide, x + 0.3, col_y + 1.4, col_w - 0.6, col_h - 1.7)
            set_p(t.text_frame.paragraphs[0], cb, tc, 15, font=font)

    add_footer(slide, page_num, title, C, font)


def build_process(prs, cfg, C, font, title, page_num):
    """Horizontal process timeline."""
    slide = new_slide(prs, C)
    steps = cfg.get("steps", [])
    if len(steps) < 2:
        return build_content(prs, cfg, C, font, title, page_num)

    if cfg.get("title"):
        t = add_tb(slide, 0.75, 0.4, 11.8, 0.6)
        set_p(t.text_frame.paragraphs[0], cfg["title"], C["text"], 30, True, font=font)

    sw, sg, sy = 2.2, 0.4, 3.85
    total_w = len(steps) * sw + (len(steps) - 1) * sg
    sx = (13.333 - total_w) / 2

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(sx), Inches(sy), Inches(total_w), Inches(0.03)
    )
    line.fill.solid(); line.fill.fore_color.rgb = C["border"]; line.line.fill.background()

    for i, step_text in enumerate(steps):
        cx = sx + i * (sw + sg)
        o = add_oval(slide, cx + sw / 2 - 0.3, sy - 0.25, 0.6, 0.6, C["primary"])
        set_p(o.text_frame.paragraphs[0], str(i + 1), C["white"], 18, True,
              PP_ALIGN.CENTER, font=font)
        o.text_frame.margin_top = Inches(0.12)
        t = add_tb(slide, cx, sy + 0.6, sw, 0.6)
        set_p(t.text_frame.paragraphs[0], step_text, C["text"], 16, False,
              PP_ALIGN.CENTER, font=font)

    note = cfg.get("note", "")
    if note:
        t = add_tb(slide, 2, sy + 1.5, 9.333, 0.4)
        set_p(t.text_frame.paragraphs[0], note, C["muted"], 11, False, PP_ALIGN.CENTER, font=font)

    add_footer(slide, page_num, title, C, font)


def build_table(prs, cfg, C, font, title, page_num):
    """Table slide."""
    slide = new_slide(prs, C)

    if cfg.get("title"):
        t = add_tb(slide, 0.75, 0.4, 11.8, 0.6)
        set_p(t.text_frame.paragraphs[0], cfg["title"], C["text"], 30, True, font=font)

    headers = cfg.get("headers", [])
    rows = cfg.get("rows", [])
    if headers and rows:
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_x, tbl_y = Inches(1.0), Inches(1.4)
        tbl_w, tbl_h = Inches(11.3), Inches(min(0.4 * n_rows + 0.3, 5.5))

        ts = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h)
        table = ts.table
        cw = int(Inches(11.3) / n_cols)
        for ci in range(n_cols):
            table.columns[ci].width = cw

        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = header
            cell.fill.solid(); cell.fill.fore_color.rgb = C["primary"]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14); p.font.bold = True
                p.font.color.rgb = C["white"]
                p.font.name = font; p.alignment = PP_ALIGN.CENTER
            cell.text_frame.margin_top = Inches(0.06)

        for ri, row in enumerate(rows):
            bg = C["panel"] if ri % 2 == 0 else C["bg"]
            for ci, cell_text in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = str(cell_text)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12); p.font.color.rgb = C["text"]
                    p.font.name = font
                    p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                cell.text_frame.margin_top = Inches(0.04)
                cell.text_frame.margin_left = Inches(0.15)

    add_footer(slide, page_num, title, C, font)


# ---- Dispatch table ----
BUILDERS = {
    "cover": build_cover,
    "section-header": build_section_header,
    "content": build_content,
    "cards": build_cards,
    "side-by-side": build_side_by_side,
    "comparison": build_comparison,
    "process": build_process,
    "table": build_table,
}


# ---- Main entry ----
def build_pptx(config: dict) -> str:
    """Build PPTX from config dict. Returns output path."""
    prs, C, font = init_prs(config)
    title = config.get("title", "Presentation")
    output = config.get("output", "output.pptx")

    for i, slide_cfg in enumerate(config.get("slides", [])):
        page_num = i + 1
        stype = slide_cfg.get("type", "content")
        builder = BUILDERS.get(stype)
        if builder:
            builder(prs, slide_cfg, C, font, title, page_num)
        else:
            print(f"WARNING: Unknown slide type '{stype}'", file=sys.stderr)

    prs.save(output)
    return output


def main():
    parser = argparse.ArgumentParser(description="Flat-function PPTX builder")
    parser.add_argument("--config", help="JSON config file path")
    parser.add_argument("--stdin", action="store_true", help="Read JSON from stdin")
    parser.add_argument("--json", action="store_true", help="Output JSON report")
    args = parser.parse_args()

    if args.stdin:
        config = json.load(sys.stdin)
    elif args.config:
        with open(args.config, "r", encoding="utf-8") as f:
            config = json.load(f)
    else:
        parser.print_help()
        sys.exit(1)

    path = build_pptx(config)
    if args.json:
        print(json.dumps({"status": "ok", "path": path,
                          "slides": len(config.get("slides", []))},
                         ensure_ascii=False))
    else:
        print(f"Generated: {path} ({len(config.get('slides', []))} slides)")


if __name__ == "__main__":
    main()
