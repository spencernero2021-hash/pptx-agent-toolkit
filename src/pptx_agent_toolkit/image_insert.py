"""
Insert images into existing PPTX slides.

Bridge between gemini:generate_image and design:create_pptx_deck.
Usage:
  python pptx_image_insert.py input.pptx output.pptx --slide 0 --image bg.jpg --layout background
  python pptx_image_insert.py input.pptx output.pptx --batch config.json
"""

import argparse
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE

# Standard widescreen (13.33 x 7.5 inches)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LAYOUTS = {}


def _register(name):
    def dec(fn):
        LAYOUTS[name] = fn
        return fn
    return dec


def _load_image(path):
    """Load image and return (width_px, height_px, content_type, image_bytes)."""
    path = Path(path)
    ext = path.suffix.lower()
    content_types = {".png": "image/png", ".jpg": "image/jpeg",
                     ".jpeg": "image/jpeg", ".gif": "image/gif",
                     ".bmp": "image/bmp", ".webp": "image/webp"}
    ct = content_types.get(ext, "image/png")
    if ext == ".webp":
        im = Image.open(path).convert("RGBA")
        import io
        buf = io.BytesIO()
        im.save(buf, format="PNG")
        return im.width, im.height, "image/png", buf.getvalue()
    with open(path, "rb") as f:
        data = f.read()
    im = Image.open(path)
    return im.width, im.height, ct, data


def _fit_rect(container_w, container_h, img_w, img_h):
    """Return (left, top, w, h) to fit image inside container, centered."""
    scale = min(container_w / img_w, container_h / img_h)
    w = int(img_w * scale)
    h = int(img_h * scale)
    left = (container_w - w) // 2
    top = (container_h - h) // 2
    return left, top, w, h


# ---- Layout presets ---------------------------------------------------------

@_register("background")
def _layout_bg(prs_w, prs_h, img_w, img_h, margin):
    """Edge-to-edge full-slide background."""
    scale = max(prs_w / img_w, prs_h / img_h)  # cover
    w = int(img_w * scale)
    h = int(img_h * scale)
    left = (prs_w - w) // 2
    top = (prs_h - h) // 2
    return left, top, w, h


@_register("left-half")
def _layout_left(prs_w, prs_h, img_w, img_h, margin):
    margin_emu = int(Inches(margin))
    region_w = prs_w // 2 - margin_emu * 2
    region_h = prs_h - margin_emu * 2
    img_left, img_top, img_w, img_h = _fit_rect(region_w, region_h, img_w, img_h)
    return margin_emu + img_left, margin_emu + img_top, img_w, img_h


@_register("right-half")
def _layout_right(prs_w, prs_h, img_w, img_h, margin):
    margin_emu = int(Inches(margin))
    half = prs_w // 2
    region_w = half - margin_emu * 2
    region_h = prs_h - margin_emu * 2
    img_left, img_top, img_w, img_h = _fit_rect(region_w, region_h, img_w, img_h)
    return half + margin_emu + img_left, margin_emu + img_top, img_w, img_h


@_register("top-half")
def _layout_top(prs_w, prs_h, img_w, img_h, margin):
    margin_emu = int(Inches(margin))
    region_w = prs_w - margin_emu * 2
    region_h = prs_h // 2 - margin_emu * 2
    img_left, img_top, img_w, img_h = _fit_rect(region_w, region_h, img_w, img_h)
    return margin_emu + img_left, margin_emu + img_top, img_w, img_h


@_register("bottom-half")
def _layout_bottom(prs_w, prs_h, img_w, img_h, margin):
    margin_emu = int(Inches(margin))
    half = prs_h // 2
    region_w = prs_w - margin_emu * 2
    region_h = half - margin_emu * 2
    img_left, img_top, img_w, img_h = _fit_rect(region_w, region_h, img_w, img_h)
    return margin_emu + img_left, half + margin_emu + img_top, img_w, img_h


@_register("center")
def _layout_center(prs_w, prs_h, img_w, img_h, margin):
    margin_emu = int(Inches(margin))
    region_w = prs_w - margin_emu * 2
    region_h = prs_h - margin_emu * 2
    img_left, img_top, img_w, img_h = _fit_rect(region_w, region_h, img_w, img_h)
    return margin_emu + img_left, margin_emu + img_top, img_w, img_h


# ---- Z-order helpers --------------------------------------------------------

def _send_to_back(slide, shape):
    """Move a shape to the bottom of the z-order (behind everything else)."""
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    # spTree children: nvGrpSpPr, grpSpPr, then shapes. Insert after grpSpPr.
    children = list(sp_tree)
    # Find index of grpSpPr (second child), insert after it
    insert_at = 2  # after nvGrpSpPr(0) and grpSpPr(1)
    sp_tree.insert(insert_at, sp)


def _send_behind_text(slide, shape):
    """Move a shape to sit behind text but above background/decorative shapes."""
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    # Insert before the last shape (typically text was added after)
    children = list(sp_tree)
    sp_tree.insert(len(children) - 1, sp)


# ---- Core functions ---------------------------------------------------------

def insert_image(prs, slide_index, image_path, layout="center", margin=0.3,
                 left=None, top=None, width=None, height=None):
    """Insert an image into a slide."""
    slide = prs.slides[slide_index]
    img_w_px, img_h_px, ct, data = _load_image(image_path)
    prs_w = prs.slide_width
    prs_h = prs.slide_height

    if layout == "inline" and all(v is not None for v in [left, top, width, height]):
        l_emu = int(Inches(left))
        t_emu = int(Inches(top))
        w_emu = int(Inches(width))
        h_emu = int(Inches(height))
    elif layout in LAYOUTS:
        l_emu, t_emu, w_emu, h_emu = LAYOUTS[layout](
            prs_w, prs_h, img_w_px, img_h_px, margin)
    else:
        raise ValueError(f"Unknown layout: {layout}. Available: {list(LAYOUTS.keys())}")

    picture = slide.shapes.add_picture(
        image_path, l_emu, t_emu, w_emu, h_emu
    )

    # Background images go behind everything; side images go behind text
    if layout == "background":
        _send_to_back(slide, picture)
    elif layout in ("left-half", "right-half", "top-half", "bottom-half"):
        _send_behind_text(slide, picture)

    return {"left_emu": l_emu, "top_emu": t_emu, "width_emu": w_emu, "height_emu": h_emu}


def process_batch(input_pptx, output_pptx, config):
    """Apply multiple image insertions from a config dict or list."""
    prs = Presentation(str(input_pptx))
    slides = config if isinstance(config, list) else config.get("slides", [])
    for item in slides:
        insert_image(
            prs,
            slide_index=item["slide"],
            image_path=item["image"],
            layout=item.get("layout", "center"),
            margin=item.get("margin", 0.3),
            left=item.get("left"),
            top=item.get("top"),
            width=item.get("width"),
            height=item.get("height"),
        )
    prs.save(str(output_pptx))
    return len(slides)


# ---- CLI --------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Insert images into PPTX slides")
    sub = parser.add_subparsers(dest="mode", required=True)

    # Single image mode
    single = sub.add_parser("single", help="Insert a single image")
    single.add_argument("input", help="Input PPTX path")
    single.add_argument("output", help="Output PPTX path")
    single.add_argument("--slide", type=int, required=True, help="Slide index (0-based)")
    single.add_argument("--image", required=True, help="Image file path")
    single.add_argument("--layout", default="center",
                        choices=list(LAYOUTS.keys()) + ["inline"],
                        help="Layout preset")
    single.add_argument("--margin", type=float, default=0.3, help="Margin in inches")
    single.add_argument("--left", type=float, help="Left in inches (inline)")
    single.add_argument("--top", type=float, help="Top in inches (inline)")
    single.add_argument("--width", type=float, help="Width in inches (inline)")
    single.add_argument("--height", type=float, help="Height in inches (inline)")

    # Batch mode
    batch = sub.add_parser("batch", help="Batch insert from JSON config")
    batch.add_argument("input", help="Input PPTX path")
    batch.add_argument("output", help="Output PPTX path")
    batch.add_argument("--config", required=True, help="JSON config file path")

    args = parser.parse_args()

    if args.mode == "single":
        prs = Presentation(args.input)
        result = insert_image(
            prs, args.slide, args.image, args.layout, args.margin,
            args.left, args.top, args.width, args.height,
        )
        prs.save(args.output)
        print(json.dumps({"status": "ok", "result": result}))

    elif args.mode == "batch":
        with open(args.config, "r", encoding="utf-8") as f:
            config = json.load(f)
        count = process_batch(args.input, args.output, config)
        print(json.dumps({"status": "ok", "slides_modified": count}))


if __name__ == "__main__":
    main()
