"""
Color contrast QA for PPTX files.

Checks every text run against its background and reports WCAG compliance.
Optionally auto-fixes low-contrast text.

Usage:
  python pptx_color_qa.py input.pptx              # Read-only report
  python pptx_color_qa.py input.pptx --json        # JSON only
  python pptx_color_qa.py input.pptx --fix out.pptx  # Auto-fix
"""

import argparse
import json
import sys
from pathlib import Path
from collections import defaultdict

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx_agent_toolkit.color_utils import (
    contrast_ratio, is_aa_pass, is_aaa_pass, find_contrast_fix, hex_to_rgb, rgb_to_hex,
)
from pptx_agent_toolkit.pptx_utils import (
    get_text_color, get_shape_fill_color, get_slide_background,
    find_bg_for_text, sample_image_behind_text, is_picture_shape,
    get_font_size, is_large_text, get_default_bg,
)


def _name(shape) -> str:
    """Get a human-readable name for a shape."""
    try:
        return shape._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr/"
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
        ).get("name", "unnamed")
    except Exception:
        return "unnamed"


def _short_text(text: str) -> str:
    """Truncate text for display."""
    if len(text) > 60:
        return text[:57] + "..."
    return text


class ColorQA:
    def __init__(self, pptx_path: str):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(self.pptx_path))
        self.default_bg = get_default_bg(self.prs)
        self.checks = []
        self.fixes = []
        self.warnings = []

    def run(self, auto_fix: bool = False, threshold: float = 4.5) -> dict:
        """Run color QA on all slides."""
        for slide_idx, slide in enumerate(self.prs.slides):
            self._check_slide(slide, slide_idx)

        if auto_fix and self.fixes:
            self._apply_fixes(threshold)

        return self._report()

    def _check_slide(self, slide, slide_idx: int):
        """Scan all text runs on a slide."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_name = _name(shape)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text.strip()
                    if not text:
                        continue
                    self._check_run(run, text, slide, slide_idx, shape, shape_name)

    def _check_run(self, run, text: str, slide, slide_idx: int, shape, shape_name: str):
        """Check a single text run for contrast."""
        fg = get_text_color(run)
        if fg is None:
            return  # Can't determine color, skip

        font_size = get_font_size(run)
        large = is_large_text(font_size)

        # Find background
        bg, bg_source = find_bg_for_text(shape, slide)

        # If background is a picture, sample image
        if bg_source == "picture" and bg is None:
            bg = sample_image_behind_text(
                slide, slide_idx, str(self.pptx_path), shape
            )
            if bg:
                bg_source = "picture_sample"

        if bg is None:
            bg = self.default_bg
            bg_source = "default"

        cr = contrast_ratio(fg, bg)
        aa = is_aa_pass(cr, large)
        aaa = is_aaa_pass(cr, large)
        level = "AAA" if aaa else ("AA" if aa else "FAIL")

        check = {
            "slide": slide_idx + 1,
            "shape": shape_name,
            "text": _short_text(text),
            "text_color": fg,
            "bg_color": bg,
            "bg_source": bg_source,
            "ratio": round(cr, 2),
            "level": level,
            "aa": aa,
            "aaa": aaa,
            "large_text": large,
            "font_size": font_size,
        }
        self.checks.append(check)

        if not aa:
            self.warnings.append(
                f"Slide {slide_idx + 1}: '{_short_text(text)}' in {shape_name} — "
                f"{fg} on {bg} = {cr:.1f}:1 (needs {'≥3.0' if large else '≥4.5'})"
            )
            self.fixes.append({
                "run": run,
                "fg": fg,
                "bg": bg,
                "text": text,
                "slide": slide_idx + 1,
                "shape": shape_name,
            })

    def _apply_fixes(self, threshold: float):
        """Apply auto-fix to all flagged runs."""
        fixed_count = 0
        for item in self.fixes:
            new_color = find_contrast_fix(item["fg"], item["bg"], threshold)
            if new_color == item["fg"]:
                continue
            run = item["run"]
            self._set_run_color(run, new_color)
            item["fixed_to"] = new_color
            fixed_count += 1

        self._fixed_count = fixed_count

    def _set_run_color(self, run, hex_color: str):
        """Set text color on a run, creating rPr if needed."""
        # Try setting at run level first for specificity
        rPr = run._r.find(f"{{{__import__('pptx_agent_toolkit.pptx_utils').NS_A}}}rPr")
        if rPr is None:
            from lxml import etree
            rPr = etree.SubElement(
                run._r,
                f"{{{__import__('pptx_agent_toolkit.pptx_utils').NS_A}}}rPr"
            )
            run._r.insert(0, rPr)

        # Remove existing fill
        for child in list(rPr):
            if child.tag.endswith("solidFill") or child.tag.endswith("ln"):
                rPr.remove(child)

        from lxml import etree
        solid = etree.SubElement(
            rPr,
            f"{{{__import__('pptx_agent_toolkit.pptx_utils').NS_A}}}solidFill"
        )
        srgb = etree.SubElement(
            solid,
            f"{{{__import__('pptx_agent_toolkit.pptx_utils').NS_A}}}srgbClr"
        )
        srgb.set("val", hex_color)

    def _report(self) -> dict:
        total = len(self.checks)
        passed = sum(1 for c in self.checks if c["aa"])
        failed = total - passed
        return {
            "status": "ok",
            "file": str(self.pptx_path),
            "total_runs": total,
            "passed": passed,
            "failed": failed,
            "threshold": "AA",
            "auto_fixed": getattr(self, "_fixed_count", 0),
            "warnings": self.warnings,
            "checks": self.checks,
        }

    def save(self, output_path: str):
        self.prs.save(str(output_path))


# ---- CLI -------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="PPTX color contrast QA")
    parser.add_argument("input", help="Input PPTX path")
    parser.add_argument("--json", action="store_true", help="JSON output only")
    parser.add_argument("--fix", help="Output path for auto-fixed PPTX")
    parser.add_argument("--threshold", type=float, default=4.5,
                        help="Contrast threshold (default 4.5 = AA normal)")
    args = parser.parse_args()

    qa = ColorQA(args.input)
    result = qa.run(auto_fix=bool(args.fix), threshold=args.threshold)

    if args.fix:
        qa.save(args.fix)
        result["fixed_path"] = args.fix

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        _print_report(result)


def _print_report(result: dict):
    """Human-readable report."""
    print(f"\n{'='*60}")
    print(f"Color Contrast QA: {Path(result['file']).name}")
    print(f"{'='*60}")
    print(f"Total text runs: {result['total_runs']}")
    print(f"Passed (AA):     {result['passed']}")
    print(f"Failed:           {result['failed']}")
    if result.get("auto_fixed"):
        print(f"Auto-fixed:       {result['auto_fixed']} runs")
        print(f"Fixed output:     {result.get('fixed_path', '')}")

    if result["warnings"]:
        print(f"\n--- Warnings ({len(result['warnings'])}) ---")
        for w in result["warnings"]:
            print(f"  ! {w}")
    else:
        print("\n  All text runs pass WCAG AA contrast. OK")

    # Summary by slide
    by_slide = defaultdict(lambda: {"total": 0, "fail": 0})
    for c in result["checks"]:
        s = c["slide"]
        by_slide[s]["total"] += 1
        if not c["aa"]:
            by_slide[s]["fail"] += 1

    print(f"\n--- By Slide ---")
    for s in sorted(by_slide):
        info = by_slide[s]
        status = "OK" if info["fail"] == 0 else f"{info['fail']} failures"
        print(f"  Slide {s}: {info['total']} runs, {status}")
    print()


if __name__ == "__main__":
    main()
