"""
PPTX Layout Enrichment Engine.

Analyzes existing PPTX slides, detects monotonous layouts, and produces an
enrichment plan with richer layout alternatives, image generation requests,
and content expansion suggestions.

Workflow:
  1. analyze()  → SlideProfile for each slide
  2. score()    → RichnessScore (0-100) per slide
  3. select()   → Best richer layout from unified catalog
  4. detect_image_needs() → Antigravity prompts for missing images
  5. expand()   → Content supplementation where needed
  6. plan()     → Full EnrichmentPlan (JSON)
  7. rebuild()  → Apply plan to generate enriched PPTX

Usage:
  python pptx_layout_enricher.py analyze input.pptx --output analysis.json
  python pptx_layout_enricher.py plan input.pptx --output plan.json
  python pptx_layout_enricher.py rebuild input.pptx --plan plan.json --output enriched.pptx
"""

import argparse
import json
import sys
import re
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
# Unified Layout Catalog
# ---------------------------------------------------------------------------
# Merged from: cookbook (16) + pptx_layout_engine (8) + design_mcp (5)
# Deduplicated on purpose/semantics. Each entry has:
#   id, name, source, purpose, content_fit, min_items, max_items,
#   needs_image, visual_weight (1-5), best_for, avoid_when
# ---------------------------------------------------------------------------

LAYOUT_CATALOG = {
    # ---- Structural ----
    "cover": {
        "id": "cover", "name": "封面", "source": "layout_engine",
        "purpose": "Opening slide with title, subtitle, accent bar",
        "content_fit": ["title", "subtitle", "eyebrow"],
        "visual_weight": 4, "needs_image": False,
        "best_for": ["演示开场", "章节封面"],
        "avoid_when": ["需要大量正文内容"],
    },
    "section-header": {
        "id": "section-header", "name": "章节分隔页", "source": "layout_engine",
        "purpose": "Centered section divider with full-width color band",
        "content_fit": ["section_title", "section_subtitle"],
        "visual_weight": 5, "needs_image": False,
        "best_for": ["章节过渡", "主题切换"],
        "avoid_when": ["需要详细内容"],
    },
    "closing": {
        "id": "closing", "name": "结束页", "source": "cookbook",
        "purpose": "Thank you / closing slide with contact info",
        "content_fit": ["thank_you", "contact"],
        "visual_weight": 4, "needs_image": False,
        "best_for": ["演示结尾", "感谢页"],
        "avoid_when": ["需要展示数据或内容"],
    },

    # ---- Content-heavy ----
    "content": {
        "id": "content", "name": "标准内容页", "source": "layout_engine",
        "purpose": "Title + bullets, classic content slide",
        "content_fit": ["title", "bullets", "body"],
        "visual_weight": 1, "needs_image": False,
        "best_for": ["纯文字列表", "需要密集信息时"],
        "avoid_when": ["可以用卡片/图文替代时优先用其他布局"],
    },
    "content-with-image": {
        "id": "content-with-image", "name": "图文内容页", "source": "layout_engine",
        "purpose": "Title + bullets on left, image on right",
        "content_fit": ["title", "bullets", "image"],
        "visual_weight": 3, "needs_image": True,
        "best_for": ["文字说明+配图", "概念解释+示意图"],
        "avoid_when": ["无合适配图时用content或cards"],
    },
    "two-column": {
        "id": "two-column", "name": "双栏对比", "source": "cookbook",
        "purpose": "Side-by-side comparison of two items/approaches",
        "content_fit": ["title", "left_content", "right_content"],
        "visual_weight": 3, "needs_image": False, "min_items": 2, "max_items": 2,
        "best_for": ["对比分析", "AB方案比较", "优劣对比"],
        "avoid_when": ["超过2个对比项", "单项内容过长"],
    },

    # ---- Visual-forward cards ----
    "cards": {
        "id": "cards", "name": "卡片页(2-3卡)", "source": "layout_engine",
        "purpose": "2-3 horizontal cards with icon + title + body",
        "content_fit": ["title", "cards"],
        "visual_weight": 3, "needs_image": False, "min_items": 2, "max_items": 3,
        "best_for": ["2-3个并列概念", "特性展示"],
        "avoid_when": ["超过3项", "卡片标题过长(>15字)"],
    },
    "multi-card": {
        "id": "multi-card", "name": "多卡片页(3-5卡)", "source": "cookbook",
        "purpose": "3-5 equal cards in a row for features/steps/pillars",
        "content_fit": ["title", "cards"],
        "visual_weight": 4, "needs_image": False, "min_items": 3, "max_items": 5,
        "best_for": ["3-5个并列要点", "特性矩阵", "服务项目"],
        "avoid_when": ["少于3项", "卡片标题>12字"],
    },
    "floating-cards": {
        "id": "floating-cards", "name": "悬浮卡片页(3卡)", "source": "cookbook",
        "purpose": "3 elevated/floating cards with depth and visual impact",
        "content_fit": ["title", "cards"],
        "visual_weight": 5, "needs_image": False, "min_items": 3, "max_items": 3,
        "best_for": ["恰好3个核心概念", "三步流程"],
        "avoid_when": ["不是恰好3项", "卡片标题>15字"],
    },

    # ---- Data & metrics ----
    "stats": {
        "id": "stats", "name": "数据指标页", "source": "cookbook",
        "purpose": "Big numbers / KPIs with labels",
        "content_fit": ["title", "metrics"],
        "visual_weight": 4, "needs_image": False, "min_items": 2, "max_items": 4,
        "best_for": ["关键指标展示", "数据亮点", "成果数据"],
        "avoid_when": ["数据需要大量上下文解释"],
    },
    "chart": {
        "id": "chart", "name": "图表页", "source": "cookbook",
        "purpose": "Data visualization with pie/doughnut/bar/column charts",
        "content_fit": ["title", "chart_data"],
        "visual_weight": 4, "needs_image": False,
        "best_for": ["数据可视化", "趋势分析", "占比展示"],
        "avoid_when": ["没有量化数据"],
    },
    "table": {
        "id": "table", "name": "表格页", "source": "layout_engine",
        "purpose": "Structured table with styled headers and alternating rows",
        "content_fit": ["title", "table_data"],
        "visual_weight": 2, "needs_image": False,
        "best_for": ["结构化数据对比", "参数列表"],
        "avoid_when": ["数据少于3行", "可以用图表更直观展示"],
    },

    # ---- Process & flow ----
    "process": {
        "id": "process", "name": "流程/时间线", "source": "layout_engine",
        "purpose": "Horizontal timeline with numbered steps",
        "content_fit": ["title", "steps"],
        "visual_weight": 3, "needs_image": False, "min_items": 2, "max_items": 6,
        "best_for": ["步骤流程", "时间线", "操作顺序"],
        "avoid_when": ["少于2步", "步骤间无先后关系"],
    },

    # ---- Dramatic / emphasis ----
    "giant-focus": {
        "id": "giant-focus", "name": "大字焦点页", "source": "cookbook",
        "purpose": "Dramatic large text for emphasis (1-3 words)",
        "content_fit": ["big_text"],
        "visual_weight": 5, "needs_image": False,
        "best_for": ["核心观点强调", "金句展示"],
        "avoid_when": ["需要详细说明"],
    },
    "bold-diagonal": {
        "id": "bold-diagonal", "name": "斜切视觉页", "source": "cookbook",
        "purpose": "High-energy diagonal shapes for urgency/impact",
        "content_fit": ["headline", "subtitle"],
        "visual_weight": 5, "needs_image": False,
        "best_for": ["关键转折点", "高能量声明"],
        "avoid_when": ["保守/传统受众"],
    },
    "corner-anchor": {
        "id": "corner-anchor", "name": "角锚不对称页", "source": "cookbook",
        "purpose": "Bold statement with strong geometric anchor in corner",
        "content_fit": ["headline", "body"],
        "visual_weight": 5, "needs_image": False,
        "best_for": ["不对称现代设计", "大胆声明"],
        "avoid_when": ["需要大量文字", "传统受众"],
    },
    "quote": {
        "id": "quote", "name": "引用页", "source": "cookbook",
        "purpose": "Pull quote with attribution and decorative elements",
        "content_fit": ["quote", "attribution"],
        "visual_weight": 3, "needs_image": False,
        "best_for": ["名人名言", "重要引述", "核心原则"],
        "avoid_when": ["不是引用内容"],
    },

    # ---- Image-forward ----
    "image-caption": {
        "id": "image-caption", "name": "图片+说明页", "source": "cookbook",
        "purpose": "Large image with caption/annotation",
        "content_fit": ["image", "caption", "title"],
        "visual_weight": 4, "needs_image": True,
        "best_for": ["产品截图", "架构图", "照片展示"],
        "avoid_when": ["无合适图片"],
    },
    "side-by-side": {
        "id": "side-by-side", "name": "图文并排页", "source": "layout_engine",
        "purpose": "Image on one side, text on the other",
        "content_fit": ["image", "bullets", "title"],
        "visual_weight": 4, "needs_image": True,
        "best_for": ["概念+示意图", "文字说明+配图"],
        "avoid_when": ["无合适图片", "文字超过6条要点"],
    },
    "comparison": {
        "id": "comparison", "name": "双面板对比页", "source": "layout_engine",
        "purpose": "Two side-by-side panels with labels (before/after, A vs B)",
        "content_fit": ["title", "left_panel", "right_panel", "labels"],
        "visual_weight": 4, "needs_image": False, "min_items": 2, "max_items": 2,
        "best_for": ["前后对比", "方案对比"],
        "avoid_when": ["只有单项内容"],
    },

    # ---- Specialized ----
    "circular-hero": {
        "id": "circular-hero", "name": "环形中心页", "source": "cookbook",
        "purpose": "Central concept with 4-6 surrounding items in circle",
        "content_fit": ["center", "surrounding_items"],
        "visual_weight": 5, "needs_image": False, "min_items": 4, "max_items": 6,
        "best_for": ["核心概念+关联要素", "生态系统展示"],
        "avoid_when": ["中心文字过长", "少于4个周边项"],
    },
    "code": {
        "id": "code", "name": "代码展示页", "source": "cookbook",
        "purpose": "Code block display with syntax-aware styling",
        "content_fit": ["code_block", "title"],
        "visual_weight": 2, "needs_image": False,
        "best_for": ["代码示例", "配置展示"],
        "avoid_when": ["非技术受众"],
    },
}

# Layout categories for quick lookup
LAYOUT_BY_CONTENT_TYPE = {
    "title_only": ["cover", "section-header", "giant-focus", "bold-diagonal", "corner-anchor"],
    "title_bullets": ["content", "content-with-image", "side-by-side", "two-column"],
    "multi_items": ["cards", "multi-card", "floating-cards"],
    "data_numbers": ["stats", "chart", "table"],
    "steps_flow": ["process"],
    "comparison": ["comparison", "two-column"],
    "image_forward": ["image-caption", "side-by-side", "content-with-image"],
    "quote_text": ["quote", "giant-focus"],
    "hub_concept": ["circular-hero"],
    "code_technical": ["code"],
}

# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class SlideProfile:
    """Analysis result for a single slide."""
    index: int
    title: str = ""
    subtitle: str = ""
    body_text: str = ""
    bullets: list = field(default_factory=list)
    shape_count: int = 0
    image_count: int = 0
    text_boxes: list = field(default_factory=list)
    has_table: bool = False
    has_chart: bool = False
    background_color: Optional[str] = None
    primary_color: Optional[str] = None
    text_color: Optional[str] = None
    font_name: str = ""
    font_sizes: list = field(default_factory=list)

    # Inferred
    current_layout_type: str = "content"
    word_count: int = 0
    is_title_slide: bool = False
    is_section_divider: bool = False
    is_closing: bool = False
    content_category: str = "title_bullets"

@dataclass
class RichnessScore:
    """Visual richness assessment for a slide."""
    index: int
    score: int  # 0-100
    level: str  # "bare", "minimal", "adequate", "rich"
    issues: list = field(default_factory=list)
    suggestions: list = field(default_factory=list)

@dataclass
class ImageNeed:
    """Image requirement for a slide."""
    slide_index: int
    priority: str  # "high", "medium", "low"
    description: str  # What the image should show
    antigravity_prompt: str  # Ready-to-use prompt for Antigravity
    suggested_layout: str  # "right-half", "left-half", "background", "center", etc.
    search_keywords: list = field(default_factory=list)  # For image library search

@dataclass
class LayoutRecommendation:
    """Recommended layout change for a slide."""
    slide_index: int
    current_layout: str
    recommended_layout: str
    reason: str
    content_mapping: dict = field(default_factory=dict)
    needs_image: bool = False
    image_need: Optional[ImageNeed] = None
    content_additions: list = field(default_factory=list)

@dataclass
class EnrichmentPlan:
    """Complete plan for enriching a PPTX."""
    source_path: str
    total_slides: int
    average_richness: float
    slides_needing_enrichment: int
    recommendations: list = field(default_factory=list)
    image_requests: list = field(default_factory=list)
    estimated_new_slides: int = 0


# ---------------------------------------------------------------------------
# Slide Analyzer
# ---------------------------------------------------------------------------

class SlideAnalyzer:
    """Extract structured content profile from PPTX slides."""

    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)

    def analyze_all(self) -> list[SlideProfile]:
        return [self.analyze_slide(i, slide) for i, slide in enumerate(self.prs.slides)]

    def analyze_slide(self, index: int, slide) -> SlideProfile:
        profile = SlideProfile(index=index)
        profile.shape_count = len(slide.shapes)

        texts = []
        for shape in slide.shapes:
            # Detect shape type
            try:
                tag = shape._element.tag.split('}')[-1] if '}' in shape._element.tag else shape._element.tag
            except Exception:
                tag = ""

            # Picture detection
            if tag == "pic":
                profile.image_count += 1
                continue

            # Table detection
            if tag == "graphicFrame":
                profile.has_table = True
                continue

            # Chart detection
            if tag == "chart":
                profile.has_chart = True
                continue

            # Text extraction
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                if full_text:
                    texts.append({
                        "text": full_text,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                    })
                    profile.text_boxes.append(full_text)

                # Font analysis
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            profile.font_sizes.append(run.font.size / 12700)  # EMU to pt
                        if run.font.name and not profile.font_name:
                            profile.font_name = run.font.name
                        if run.font.color and not profile.text_color:
                            try:
                                if run.font.color.rgb:
                                    profile.text_color = str(run.font.color.rgb)
                            except (AttributeError, TypeError, ValueError):
                                pass

        # Classify texts
        if texts:
            # Title: typically top-most, largest text
            sorted_by_top = sorted(texts, key=lambda t: t["top"])
            largest = max(texts, key=lambda t: t["height"] if t["height"] else 0)
            profile.title = sorted_by_top[0]["text"] if sorted_by_top else ""

            # Body/secondary text
            if len(texts) > 1:
                profile.body_text = texts[1]["text"] if len(texts) > 1 else ""

            # Detect bullets (lines starting with bullet markers or short lines)
            for tb in texts:
                lines = tb["text"].split('\n')
                for line in lines:
                    stripped = line.strip()
                    if stripped and len(stripped) > 3:
                        profile.bullets.append(stripped)

        # Word count
        all_text = ' '.join(tb["text"] for tb in texts)
        profile.word_count = len(all_text)

        # Infer current layout type
        profile.current_layout_type = self._infer_layout(profile)
        profile.content_category = self._infer_category(profile)
        profile.is_title_slide = (index == 0 and profile.image_count == 0 and profile.word_count < 30)
        profile.is_section_divider = (profile.word_count < 20 and profile.shape_count <= 3 and not profile.bullets)
        profile.is_closing = (index > 0 and any(
            kw in (profile.title + profile.body_text).lower()
            for kw in ["谢谢", "thank", "感谢", "q&a", "问答", "结束"]
        ))

        return profile

    def _infer_layout(self, p: SlideProfile) -> str:
        """Infer current layout type from slide structure."""
        if p.is_title_slide:
            return "cover"
        if p.is_section_divider:
            return "section-header"
        if p.has_table:
            return "table"
        if p.has_chart:
            return "chart"
        if p.image_count >= 1 and len(p.bullets) >= 2:
            return "side-by-side"
        if p.image_count >= 1:
            return "image-caption"
        if len(p.bullets) >= 5:
            return "content"
        if len(p.bullets) >= 3 and p.shape_count <= 3:
            return "content"
        return "content"

    def _infer_category(self, p: SlideProfile) -> str:
        if len(p.bullets) >= 4:
            return "title_bullets"
        if p.image_count > 0:
            return "image_forward"
        return "title_bullets"


# ---------------------------------------------------------------------------
# Richness Scorer
# ---------------------------------------------------------------------------

class RichnessScorer:
    """Score slides 0-100 on visual richness."""

    def score(self, profile: SlideProfile) -> RichnessScore:
        score = 50  # baseline
        issues = []
        suggestions = []

        # + Visual elements
        if profile.image_count > 0:
            score += 20
        if profile.has_table or profile.has_chart:
            score += 15
        if profile.shape_count > 3:
            score += min((profile.shape_count - 3) * 3, 15)

        # + Layout variety (relative to content type)
        if profile.current_layout_type not in ("content",):
            score += 10

        # - Penalties
        if profile.image_count == 0 and profile.shape_count <= 3 and len(profile.bullets) >= 3:
            score -= 25
            issues.append("纯文字页，缺少视觉元素")
            suggestions.append("添加配图或改用卡片布局")

        if profile.image_count == 0:
            if len(profile.bullets) >= 4:
                issues.append("多条要点无配图")
                suggestions.append("将要点拆分为卡片或用图文并排布局")

        if profile.shape_count <= 2 and len(profile.bullets) > 0:
            score -= 10
            issues.append("形状数量极少，排版单一")
            suggestions.append("增加装饰性几何元素或使用更丰富的布局")

        if profile.word_count > 200:
            score -= 10
            issues.append("文字密度过高")
            suggestions.append("精简文字或将内容拆分到多页")

        if profile.word_count < 5 and not profile.is_section_divider and not profile.is_title_slide:
            score -= 5
            issues.append("内容过少")
            suggestions.append("补充说明文字或相关配图")

        # Clamp
        score = max(0, min(100, score))

        # Level
        if score < 30:
            level = "bare"
        elif score < 50:
            level = "minimal"
        elif score < 75:
            level = "adequate"
        else:
            level = "rich"

        return RichnessScore(
            index=profile.index,
            score=score,
            level=level,
            issues=issues,
            suggestions=suggestions,
        )

    def score_all(self, profiles: list[SlideProfile]) -> list[RichnessScore]:
        return [self.score(p) for p in profiles]


# ---------------------------------------------------------------------------
# Layout Selector
# ---------------------------------------------------------------------------

class LayoutSelector:
    """Select the best richer layout for slide content."""

    # Priority: richer layouts considered first
    ENRICHMENT_PRIORITY = [
        # Visual-forward (prefer these for text-only content)
        "floating-cards", "multi-card", "cards",
        "side-by-side", "content-with-image",
        "stats", "circular-hero",
        # Dramatic
        "giant-focus", "bold-diagonal", "corner-anchor", "quote",
        # Structured
        "comparison", "two-column", "process",
        "image-caption",
        # Data
        "chart", "table",
        # Fallback
        "content",
    ]

    # Mapping: content pattern → best layout IDs
    CONTENT_PATTERN_MAP = {
        # Multiple equal items
        "3_equal_items": ["floating-cards", "cards"],
        "4_5_equal_items": ["multi-card", "cards"],
        "2_equal_items": ["comparison", "two-column"],
        # Bullets list
        "many_bullets": ["side-by-side", "content-with-image", "multi-card", "content"],
        "few_bullets": ["giant-focus", "corner-anchor", "side-by-side", "content"],
        # Numbers/Data
        "kpi_numbers": ["stats", "chart"],
        "tabular_data": ["table", "chart"],
        # Process
        "sequential_steps": ["process", "floating-cards"],
        # Hub
        "hub_concept": ["circular-hero", "multi-card"],
        # Quotes
        "quote_text": ["quote", "giant-focus", "corner-anchor"],
        # Image needed
        "with_image": ["image-caption", "side-by-side", "content-with-image"],
    }

    def select(self, profile: SlideProfile) -> LayoutRecommendation:
        current = profile.current_layout_type

        # Determine content pattern
        pattern = self._detect_pattern(profile)

        # Get candidate layouts
        candidates = self.CONTENT_PATTERN_MAP.get(pattern, ["content"])

        # Filter by applicability
        applicable = []
        for cid in candidates:
            if cid not in LAYOUT_CATALOG:
                continue
            layout = LAYOUT_CATALOG[cid]
            # Check min/max item constraints
            if "min_items" in layout and len(profile.bullets) < layout["min_items"]:
                continue
            if "max_items" in layout and len(profile.bullets) > layout["max_items"]:
                continue
            # Skip if needs image but we'd need to detect that separately
            applicable.append(cid)

        if not applicable:
            applicable = ["content"]

        # Pick the highest-priority applicable layout
        chosen = applicable[0]
        for cid in self.ENRICHMENT_PRIORITY:
            if cid in applicable:
                chosen = cid
                break

        layout = LAYOUT_CATALOG.get(chosen, LAYOUT_CATALOG["content"])

        # Determine if image is needed
        needs_image = layout.get("needs_image", False)

        # Generate reason
        if chosen == current:
            reason = f"当前已是{layout['name']}，保持不变"
        else:
            reason = f"从{LAYOUT_CATALOG.get(current, {}).get('name', current)}升级为{layout['name']}: {layout['purpose']}"

        # Build content mapping
        content_map = self._build_content_map(profile, chosen)

        return LayoutRecommendation(
            slide_index=profile.index,
            current_layout=current,
            recommended_layout=chosen,
            reason=reason,
            content_mapping=content_map,
            needs_image=needs_image,
        )

    def _detect_pattern(self, profile: SlideProfile) -> str:
        n = len(profile.bullets)
        if n == 3:
            return "3_equal_items"
        if 4 <= n <= 5:
            return "4_5_equal_items"
        if n == 2:
            return "2_equal_items"
        if n > 5:
            return "many_bullets"
        if n == 0:
            return "few_bullets"
        return "many_bullets"

    def _build_content_map(self, profile: SlideProfile, layout_id: str) -> dict:
        """Map profile content to layout's content_fit fields."""
        layout = LAYOUT_CATALOG.get(layout_id, {})
        fits = layout.get("content_fit", [])
        mapping = {}

        for fit in fits:
            if fit == "title":
                mapping["title"] = profile.title
            elif fit == "bullets" and profile.bullets:
                mapping["bullets"] = profile.bullets[:8]  # cap at 8
            elif fit == "body":
                mapping["body"] = profile.body_text or profile.title
            elif fit == "cards" and profile.bullets:
                cards = []
                for i, b in enumerate(profile.bullets[:5]):
                    # Split bullet into title (first ~10 chars) and body
                    if len(b) > 15:
                        cards.append({"title": b[:12], "body": b, "icon": f"{i+1:02d}"})
                    else:
                        cards.append({"title": b, "body": "", "icon": f"{i+1:02d}"})
                mapping["cards"] = cards
            elif fit == "steps" and profile.bullets:
                mapping["steps"] = profile.bullets[:6]
            elif fit == "metrics" and profile.bullets:
                metrics = []
                for b in profile.bullets[:4]:
                    # Try to extract number from bullet
                    match = re.search(r'(\d+[\d,.]*\s*[%x倍万千亿]?)', b)
                    if match:
                        metrics.append({"value": match.group(1), "label": b[:30]})
                if metrics:
                    mapping["metrics"] = metrics

        return mapping


# ---------------------------------------------------------------------------
# Image Needs Detector
# ---------------------------------------------------------------------------

class ImageNeedsDetector:
    """Detect slides that would benefit from images and generate Antigravity prompts."""

    # Keywords that suggest what kind of image is needed
    TOPIC_TO_IMAGE_PROMPT = {
        "实验": "laboratory experiment setup, scientific equipment, clean professional lab photography",
        "细胞": "microscopic cell structure, biology cellular diagram, scientific illustration",
        "基因": "DNA helix structure, genetic science visualization, molecular biology",
        "数据": "data visualization dashboard, analytics charts, modern data science",
        "流程": "workflow process diagram, step-by-step flowchart, business process illustration",
        "结构": "architectural structure diagram, system architecture, organizational chart",
        "技术": "modern technology abstract, circuit board or digital network visualization",
        "安全": "cybersecurity shield, data protection concept, security system diagram",
        "对比": "side-by-side comparison illustration, before-after concept visual",
        "方法": "methodology framework diagram, scientific method illustration",
    }

    def detect(self, profile: SlideProfile, recommendation: LayoutRecommendation) -> list[ImageNeed]:
        needs = []

        # Already has images
        if profile.image_count > 0:
            return needs

        # Only suggest images for layouts that benefit from them
        if not recommendation.needs_image and profile.image_count == 0:
            # Still check: text-heavy slides benefit from images
            if len(profile.bullets) >= 4 and recommendation.recommended_layout in ("content", "content-with-image"):
                recommendation.needs_image = True

        if not recommendation.needs_image:
            return needs

        # Generate image prompt from slide content
        prompt = self._generate_prompt(profile)
        if not prompt:
            return needs

        need = ImageNeed(
            slide_index=profile.index,
            priority="high" if profile.word_count > 100 else "medium",
            description=f"为\"{profile.title[:30]}\"页面生成配图",
            antigravity_prompt=prompt,
            suggested_layout="right-half",
            search_keywords=self._extract_keywords(profile),
        )
        return [need]

    def _generate_prompt(self, profile: SlideProfile) -> str:
        """Generate an Antigravity-compatible image prompt from slide content."""
        title = profile.title
        body = profile.body_text[:100] if profile.body_text else ""

        # Match topic keywords
        for keyword, prompt_template in self.TOPIC_TO_IMAGE_PROMPT.items():
            if keyword in title or keyword in body:
                return (
                    f"Professional presentation illustration: {prompt_template}. "
                    f"Clean background, modern style, suitable for slide background. "
                    f"High quality, well-lit, composition with copy space on the side."
                )

        # Generic fallback
        return (
            f"Professional abstract illustration related to '{title}'. "
            f"Modern corporate style, clean composition, subtle colors, "
            f"suitable for presentation slide background with copy space. High quality."
        )

    def _extract_keywords(self, profile: SlideProfile) -> list[str]:
        """Extract search keywords from slide content."""
        keywords = []
        title_words = re.findall(r'[一-鿿\w]+', profile.title)
        keywords.extend(title_words[:5])
        if profile.body_text:
            body_words = re.findall(r'[一-鿿]{2,}', profile.body_text)
            keywords.extend(body_words[:3])
        return keywords[:6]


# ---------------------------------------------------------------------------
# Content Expander
# ---------------------------------------------------------------------------

class ContentExpander:
    """Generate supplementary content for sparse slides."""

    def expand(self, profile: SlideProfile, recommendation: LayoutRecommendation) -> list[str]:
        additions = []

        # Sparse content: add subtitle if missing
        if profile.word_count < 15 and not profile.is_title_slide and not profile.is_section_divider:
            if not profile.subtitle and profile.title:
                additions.append(f"subtitle: {profile.title}的关键要点与详细说明")

        # Missing body for title+bullets slides
        if len(profile.bullets) >= 3 and not profile.body_text:
            additions.append(f"body: 以下内容将从多个维度对{profile.title}进行阐述")

        # Content expansion for card layouts
        if recommendation.recommended_layout in ("cards", "multi-card", "floating-cards"):
            if len(profile.bullets) < 3:
                additions.append("expand_items: 将内容拆分为3-4个卡片，每个卡片包含标题和简短描述")

        return additions


# ---------------------------------------------------------------------------
# Enrichment Engine (orchestrator)
# ---------------------------------------------------------------------------

class EnrichmentEngine:
    """Orchestrates analysis → scoring → recommendations → plan."""

    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.analyzer = SlideAnalyzer(pptx_path)
        self.scorer = RichnessScorer()
        self.selector = LayoutSelector()
        self.image_detector = ImageNeedsDetector()
        self.expander = ContentExpander()

    def analyze(self) -> list[SlideProfile]:
        return self.analyzer.analyze_all()

    def plan(self) -> EnrichmentPlan:
        profiles = self.analyze()
        scores = self.scorer.score_all(profiles)
        recommendations = []
        image_requests = []

        for profile in profiles:
            score = scores[profile.index]
            rec = self.selector.select(profile)

            # Only recommend changes for non-rich slides
            if score.level in ("bare", "minimal") and rec.recommended_layout != profile.current_layout_type:
                # Detect image needs
                needs = self.image_detector.detect(profile, rec)
                if needs:
                    rec.image_need = needs[0]
                    image_requests.extend(needs)

                # Expand content
                additions = self.expander.expand(profile, rec)
                rec.content_additions = additions

                recommendations.append(rec)

        avg_richness = sum(s.score for s in scores) / max(len(scores), 1)

        # Estimate new slides (some sparse slides split into 2)
        new_slides = len(profiles)
        for rec in recommendations:
            if "拆分" in ' '.join(rec.content_additions):
                new_slides += 1

        return EnrichmentPlan(
            source_path=self.pptx_path,
            total_slides=len(profiles),
            average_richness=round(avg_richness, 1),
            slides_needing_enrichment=len(recommendations),
            recommendations=recommendations,
            image_requests=image_requests,
            estimated_new_slides=new_slides,
        )

    def plan_to_dict(self, plan: EnrichmentPlan) -> dict:
        """Convert plan to JSON-serializable dict."""
        return {
            "source": plan.source_path,
            "total_slides": plan.total_slides,
            "average_richness": plan.average_richness,
            "slides_needing_enrichment": plan.slides_needing_enrichment,
            "estimated_new_slides": plan.estimated_new_slides,
            "recommendations": [
                {
                    "slide_index": r.slide_index,
                    "current_layout": r.current_layout,
                    "recommended_layout": r.recommended_layout,
                    "layout_name": LAYOUT_CATALOG.get(r.recommended_layout, {}).get("name", r.recommended_layout),
                    "reason": r.reason,
                    "content_mapping": r.content_mapping,
                    "needs_image": r.needs_image,
                    "image_need": {
                        "description": r.image_need.description,
                        "antigravity_prompt": r.image_need.antigravity_prompt,
                        "priority": r.image_need.priority,
                        "suggested_layout": r.image_need.suggested_layout,
                    } if r.image_need else None,
                    "content_additions": r.content_additions,
                }
                for r in plan.recommendations
            ],
            "image_requests": [
                {
                    "slide_index": img.slide_index,
                    "description": img.description,
                    "antigravity_prompt": img.antigravity_prompt,
                    "priority": img.priority,
                    "suggested_layout": img.suggested_layout,
                }
                for img in plan.image_requests
            ],
        }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def cmd_analyze(args):
    """Analyze PPTX and output slide profiles."""
    engine = EnrichmentEngine(args.input)
    profiles = engine.analyze()
    scores = engine.scorer.score_all(profiles)

    output = {
        "source": args.input,
        "slide_count": len(profiles),
        "slides": [],
    }

    for p, s in zip(profiles, scores):
        output["slides"].append({
            "index": p.index,
            "title": p.title,
            "word_count": p.word_count,
            "shape_count": p.shape_count,
            "image_count": p.image_count,
            "bullet_count": len(p.bullets),
            "current_layout": p.current_layout_type,
            "richness_score": s.score,
            "richness_level": s.level,
            "issues": s.issues,
        })

    out_path = args.output or "analysis.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)
    print(f"Analysis saved to {out_path}")
    print(f"Slides: {len(profiles)}, Avg richness: {sum(s.score for s in scores)/len(scores):.0f}")


def cmd_plan(args):
    """Generate enrichment plan."""
    engine = EnrichmentEngine(args.input)
    plan = engine.plan()
    plan_dict = engine.plan_to_dict(plan)

    out_path = args.output or "enrichment_plan.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(plan_dict, f, ensure_ascii=False, indent=2)
    print(f"Plan saved to {out_path}")
    print(f"Total slides: {plan.total_slides}")
    print(f"Average richness: {plan.average_richness}/100")
    print(f"Slides needing enrichment: {plan.slides_needing_enrichment}")
    print(f"Image requests: {len(plan.image_requests)}")

    if plan.recommendations:
        print("\nTop recommendations:")
        for r in plan.recommendations[:5]:
            name = LAYOUT_CATALOG.get(r.recommended_layout, {}).get("name", r.recommended_layout)
            print(f"  Slide {r.slide_index+1}: {r.current_layout} → {name}")
            print(f"    {r.reason}")
            if r.image_need:
                print(f"    [IMAGE] {r.image_need.description}")


def cmd_list_layouts(args):
    """List all available layouts."""
    filter_source = args.source or ""
    print(f"{'ID':<22} {'Name':<12} {'Source':<16} {'Weight':<8} {'NeedsImg'}")
    print("-" * 78)
    for lid, layout in LAYOUT_CATALOG.items():
        if filter_source and layout["source"] != filter_source:
            continue
        print(f"{lid:<22} {layout['name']:<12} {layout['source']:<16} "
              f"{layout['visual_weight']:<8} {layout['needs_image']}")


def main():
    parser = argparse.ArgumentParser(description="PPTX Layout Enrichment Engine")
    sub = parser.add_subparsers(dest="command")

    # analyze
    p_analyze = sub.add_parser("analyze", help="Analyze PPTX structure")
    p_analyze.add_argument("input", help="Input PPTX path")
    p_analyze.add_argument("--output", "-o", help="Output JSON path")

    # plan
    p_plan = sub.add_parser("plan", help="Generate enrichment plan")
    p_plan.add_argument("input", help="Input PPTX path")
    p_plan.add_argument("--output", "-o", help="Output JSON path")

    # list-layouts
    p_list = sub.add_parser("list-layouts", help="List available layouts")
    p_list.add_argument("--source", "-s", help="Filter by source (cookbook, layout_engine, design_mcp)")

    args = parser.parse_args()

    if args.command == "analyze":
        cmd_analyze(args)
    elif args.command == "plan":
        cmd_plan(args)
    elif args.command == "list-layouts":
        cmd_list_layouts(args)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
