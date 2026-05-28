"""
Rich-layout PPTX generator with WCAG-safe color palettes.

8 layout types, 5 palettes, generates standalone PPTX from JSON config.

Usage:
  python pptx_layout_engine.py --config layout.json
  python pptx_layout_engine.py --stdin < layout.json
"""

import argparse
import json
import sys
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx_agent_toolkit.color_utils import palette_aa_check, contrast_ratio, hex_to_rgb

# ---- Palettes (all WCAG-validated) ----------------------------------------

PALETTES = {
    "premium": {
        "text": "161A1D",       # near-black on light bg
        "muted": "596168",      # gray
        "bg": "F6F3EC",         # warm cream
        "panel": "FFFCF5",      # off-white card
        "primary": "1A3A5C",    # deep navy
        "accent": "D66A2E",     # warm orange
        "accent2": "245C4D",    # dark teal
        "text_on_primary": "FFFFFF",
        "text_on_accent": "FFFFFF",
        "border": "E6DED1",     # subtle tan
        "font_name": "Microsoft YaHei",
    },
    "modern": {
        "text": "1A1A2E",       # near-black
        "muted": "6B7280",      # cool gray
        "bg": "FFFFFF",         # pure white
        "panel": "F8F9FA",      # light gray card
        "primary": "2563EB",    # electric blue
        "accent": "7C3AED",     # purple
        "accent2": "0EA5E9",    # sky blue
        "text_on_primary": "FFFFFF",
        "text_on_accent": "FFFFFF",
        "border": "E5E7EB",
        "font_name": "Microsoft YaHei",
    },
    "warm": {
        "text": "1C1917",       # warm black
        "muted": "6B5E56",      # warm gray (adjusted for AA)
        "bg": "FAF7F2",         # warm white
        "panel": "FFFBF5",      # warm card
        "primary": "78350F",    # dark amber
        "accent": "C2410C",     # orange-red
        "accent2": "92400E",    # dark orange
        "text_on_primary": "FFFFFF",
        "text_on_accent": "FFFFFF",
        "border": "E7E0D5",
        "font_name": "Microsoft YaHei",
    },
    "dark": {
        "text": "F0F0F0",       # light text on dark bg
        "muted": "9CA3AF",      # muted gray
        "bg": "111827",         # dark navy background
        "panel": "1F2937",      # slightly lighter panel
        "primary": "1D4ED8",    # blue (dark enough for white text AA)
        "accent": "F59E0B",     # amber
        "accent2": "10B981",    # emerald
        "text_on_primary": "FFFFFF",
        "text_on_accent": "111827",
        "border": "374151",
        "font_name": "Microsoft YaHei",
    },
    "corporate": {
        "text": "0F172A",       # slate dark
        "muted": "475569",      # slate gray
        "bg": "FFFFFF",         # white
        "panel": "F1F5F9",      # slate light
        "primary": "1E40AF",    # corporate blue
        "accent": "0D9488",     # teal
        "accent2": "6366F1",    # indigo
        "text_on_primary": "FFFFFF",
        "text_on_accent": "FFFFFF",
        "border": "CBD5E1",
        "font_name": "Microsoft YaHei",
    },
}


def validate_palette(name: str, p: dict) -> list[str]:
    """Check that critical pairs in a palette meet AA contrast."""
    test = {
        "text": p["text"],
        "bg": p["bg"],
        "muted": p["muted"],
        "primary": p["primary"],
        "panel": p["panel"],
        "text_on_primary": p["text_on_primary"],
    }
    return palette_aa_check(test)


# Validate all built-in palettes at import time
_palette_warnings = {}
for _name, _pal in PALETTES.items():
    _w = validate_palette(_name, _pal)
    if _w:
        _palette_warnings[_name] = _w
if _palette_warnings:
    print("WARNING: Palette contrast issues:", file=sys.stderr)
    for _n, _ws in _palette_warnings.items():
        for _w in _ws:
            print(f"  {_n}: {_w}", file=sys.stderr)


# ---- Layout Engine ---------------------------------------------------------

class LayoutEngine:
    def __init__(self, config: dict):
        self.config = config
        palette_name = config.get("palette", "premium")
        self.palette = dict(PALETTES.get(palette_name, PALETTES["premium"]))

        # Apply overrides
        overrides = config.get("palette_overrides", {})
        self.palette.update(overrides)

        # Re-validate if overridden
        if overrides:
            w = validate_palette("custom", self.palette)
            for msg in w:
                print(f"WARNING: {msg}", file=sys.stderr)

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.output = config.get("output", "output.pptx")
        self.title = config.get("title", "")
        self.page = 0

    def build(self) -> str:
        """Generate all slides. Returns output path."""
        for slide_conf in self.config.get("slides", []):
            self.page += 1
            stype = slide_conf.get("type", "content")
            method = getattr(self, f"_slide_{stype.replace('-', '_')}", None)
            if method:
                method(slide_conf)
            else:
                print(f"WARNING: Unknown slide type '{stype}', using content",
                      file=sys.stderr)
                self._slide_content(slide_conf)

        self.prs.save(self.output)
        return self.output

    def report(self) -> dict:
        """Return summary after build."""
        return {
            "status": "ok",
            "path": self.output,
            "slides": self.page,
            "palette": self.config.get("palette", "premium"),
        }

    # ---- Helpers -----------------------------------------------------------

    def _hex(self, key: str) -> str:
        return self.palette.get(key, "000000")

    def _rgb(self, key: str) -> RGBColor:
        h = self._hex(key)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _blank_slide(self) -> object:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        # Draw a full-slide background rectangle (avoids XML manipulation)
        bg_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height,
        )
        self._shape_fill(bg_rect, "bg")
        self._shape_line(bg_rect, None)
        # Send to back of z-order
        sp_tree = slide.shapes._spTree
        sp = bg_rect._element
        sp_tree.remove(sp)
        sp_tree.insert(2, sp)  # after nvGrpSpPr(0) and grpSpPr(1)
        return slide

    def _add_footer(self, slide):
        """Page number + title footer."""
        txt = f"{self.title} / {self.page:02d}"
        cx = Inches(11.8) if self.page < 10 else Inches(1.5)
        tb = self._add_text_box(
            slide, None,
            Inches(0.75), Inches(7.05), cx, Inches(0.25),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(10)
        p.font.color.rgb = self._rgb("muted")
        p.font.name = self.palette.get("font_name", "Microsoft YaHei")

    def _add_text_box(self, slide, fill_key: str | None,
                       left, top, width, height) -> object:
        """Add a text box with optional background fill."""
        left_emu = Inches(left) if isinstance(left, (int, float)) else left
        top_emu = Inches(top) if isinstance(top, (int, float)) else top
        w_emu = Inches(width) if isinstance(width, (int, float)) else width
        h_emu = Inches(height) if isinstance(height, (int, float)) else height

        tb = slide.shapes.add_textbox(left_emu, top_emu, w_emu, h_emu)
        if fill_key:
            self._shape_fill(tb, fill_key)
        tb.text_frame.word_wrap = True
        tb.text_frame.margin_left = Inches(0.15)
        tb.text_frame.margin_right = Inches(0.15)
        tb.text_frame.margin_top = Inches(0.1)
        tb.text_frame.margin_bottom = Inches(0.1)
        return tb

    def _shape_fill(self, shape, color_key: str):
        """Set solid fill on a shape by palette key."""
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(color_key)
        except Exception:
            pass

    def _shape_line(self, shape, color_key: str | None, width_pt: float = 1):
        """Set or clear shape outline."""
        try:
            if color_key is None:
                shape.line.fill.background()
            else:
                shape.line.color.rgb = self._rgb(color_key)
                shape.line.width = Pt(width_pt)
        except Exception:
            pass

    def _set_text(self, text_frame, text: str, color_key: str = "text",
                  size_pt: float = 15, bold: bool = False,
                  alignment=PP_ALIGN.LEFT, spacing_pt: float | None = None):
        """Set text on a paragraph with full formatting."""
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size_pt)
        p.font.bold = bold
        p.font.color.rgb = self._rgb(color_key)
        p.font.name = self.palette.get("font_name", "Microsoft YaHei")
        p.alignment = alignment
        if spacing_pt is not None:
            p.space_after = Pt(spacing_pt)
        return p

    def _add_para(self, text_frame, text: str, color_key: str = "text",
                  size_pt: float = 14, bold: bool = False,
                  alignment=PP_ALIGN.LEFT, spacing_pt: float = 4):
        """Add a new paragraph."""
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(size_pt)
        p.font.bold = bold
        p.font.color.rgb = self._rgb(color_key)
        p.font.name = self.palette.get("font_name", "Microsoft YaHei")
        p.alignment = alignment
        if spacing_pt is not None:
            p.space_after = Pt(spacing_pt)
        return p

    # ---- Slide layouts -----------------------------------------------------

    def _slide_cover(self, conf):
        """Full-width cover slide with accent bar."""
        slide = self._blank_slide()
        # Full background band on left
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, Inches(0.12), Inches(7.5),
        )
        self._shape_fill(band, "primary")
        self._shape_line(band, None)

        # Eyebrow
        eyebrow = conf.get("eyebrow", "")
        if eyebrow:
            tb = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(2.0), Inches(11.8), Inches(0.4),
            )
            self._set_text(tb.text_frame, eyebrow, "primary", 13, True,
                           PP_ALIGN.LEFT)

        # Title
        tb = self._add_text_box(
            slide, None,
            Inches(0.75), Inches(2.5), Inches(11.8), Inches(1.2),
        )
        title = conf.get("title", self.title)
        self._set_text(tb.text_frame, title, "text", 42, True, PP_ALIGN.LEFT)

        # Subtitle
        subtitle = conf.get("subtitle", "")
        if subtitle:
            tb2 = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(3.8), Inches(11.8), Inches(0.5),
            )
            self._set_text(tb2.text_frame, subtitle, "muted", 18, False,
                           PP_ALIGN.LEFT)

        # Accent line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(4.5), Inches(1.2), Inches(0.05),
        )
        self._shape_fill(line, "accent")
        self._shape_line(line, None)

        # Kicker / slide count
        kicker = conf.get("kicker", "")
        if kicker:
            tb3 = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(5.2), Inches(3), Inches(0.3),
            )
            self._set_text(tb3.text_frame, kicker, "muted", 12, False,
                           PP_ALIGN.LEFT)

        self._add_footer(slide)

    def _slide_section_header(self, conf):
        """Centered section divider."""
        slide = self._blank_slide()
        # Full-width colored band
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(2.2), Inches(13.333), Inches(3.1),
        )
        self._shape_fill(band, "primary")
        self._shape_line(band, None)

        # Title centered
        tb = self._add_text_box(
            slide, "primary",
            Inches(1.5), Inches(2.7), Inches(10.333), Inches(1.0),
        )
        self._set_text(tb.text_frame, conf.get("title", ""),
                       "text_on_primary", 40, True, PP_ALIGN.CENTER)

        # Subtitle
        subtitle = conf.get("subtitle", "")
        if subtitle:
            tb2 = self._add_text_box(
                slide, "primary",
                Inches(1.5), Inches(3.8), Inches(10.333), Inches(0.5),
            )
            self._set_text(tb2.text_frame, subtitle, "text_on_primary",
                           16, False, PP_ALIGN.CENTER)

        self._add_footer(slide)

    def _slide_content(self, conf):
        """Title + bullets (classic)."""
        slide = self._blank_slide()
        page_num = conf.get("number", self.page)

        # Section number
        tb = self._add_text_box(
            slide, None,
            Inches(0.75), Inches(0.4), Inches(1), Inches(0.4),
        )
        num_str = f"{page_num:02d}" if isinstance(page_num, int) else str(page_num)
        self._set_text(tb.text_frame, num_str, "primary", 13, True)

        # Title
        tb2 = self._add_text_box(
            slide, None,
            Inches(0.75), Inches(0.9), Inches(11.8), Inches(0.7),
        )
        self._set_text(tb2.text_frame, conf.get("title", ""), "text", 34, True)

        # Accent bar under title
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(1.6), Inches(1.2), Inches(0.04),
        )
        self._shape_fill(bar, "accent")
        self._shape_line(bar, None)

        # Body text
        body = conf.get("body", "")
        y = 1.9
        if body:
            tb3 = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(y), Inches(11.8), Inches(0.5),
            )
            self._set_text(tb3.text_frame, body, "muted", 15, False)
            y += 0.6

        # Bullets
        bullets = conf.get("bullets", [])
        if bullets:
            tb4 = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(y), Inches(11.8), Inches(4.5 - (y - 1.9)),
            )
            first = True
            for b in bullets:
                if first:
                    self._set_text(tb4.text_frame, f"  {b}", "text", 18, False,
                                   PP_ALIGN.LEFT, spacing_pt=12)
                    first = False
                else:
                    self._add_para(tb4.text_frame, f"  {b}", "text", 18, False,
                                   spacing_pt=12)

        self._add_footer(slide)

    def _slide_cards(self, conf):
        """2-3 horizontal cards with icon + title + body."""
        slide = self._blank_slide()
        items = conf.get("items", [])
        n = len(items)
        if n not in (2, 3):
            n = 3

        # Title
        if conf.get("title"):
            tb = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.7),
            )
            self._set_text(tb.text_frame, conf["title"], "text", 32, True)

        # Card layout
        margin = 0.75
        gap = 0.3
        total_w = 11.8
        card_w = (total_w - gap * (n - 1)) / n
        card_h = 4.8
        card_y = 1.3

        for i, item in enumerate(items):
            x = margin + i * (card_w + gap)

            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(card_y),
                Inches(card_w), Inches(card_h),
            )
            self._shape_fill(card, "panel")
            self._shape_line(card, "border", 0.5)

            # Icon circle (number)
            icon = item.get("icon", f"{i + 1:02d}")
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.3), Inches(card_y + 0.3),
                Inches(0.6), Inches(0.6),
            )
            self._shape_fill(circle, "primary")
            self._shape_line(circle, None)
            self._set_text(circle.text_frame, icon, "text_on_primary",
                           18, True, PP_ALIGN.CENTER)
            circle.text_frame.margin_top = Inches(0.12)

            # Card title
            tb2 = self._add_text_box(
                slide, None,
                Inches(x + 0.3), Inches(card_y + 1.1),
                Inches(card_w - 0.6), Inches(0.4),
            )
            self._set_text(tb2.text_frame, item.get("title", ""),
                           "text", 18, True)

            # Card body
            body = item.get("body", "")
            if body:
                tb3 = self._add_text_box(
                    slide, None,
                    Inches(x + 0.3), Inches(card_y + 1.6),
                    Inches(card_w - 0.6), Inches(card_h - 1.9),
                )
                self._set_text(tb3.text_frame, body, "muted", 14, False)

            # Accent line bottom
            aline = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + card_w - 1.5), Inches(card_y + card_h - 0.06),
                Inches(1.2), Inches(0.04),
            )
            self._shape_fill(aline, "accent")
            self._shape_line(aline, None)

        self._add_footer(slide)

    def _slide_side_by_side(self, conf):
        """Image on one side, text on the other."""
        slide = self._blank_slide()
        image_side = conf.get("image_side", "right")
        image_path = conf.get("image", "")

        # Title
        if conf.get("title"):
            tb = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.6),
            )
            self._set_text(tb.text_frame, conf["title"], "text", 30, True)

        # Image side
        img_w = 5.5
        img_h = 4.5
        img_x = 0.75 if image_side == "left" else (13.333 - 0.75 - img_w)
        img_y = 1.4

        if image_path and Path(image_path).exists():
            try:
                pic = slide.shapes.add_picture(
                    str(image_path),
                    Inches(img_x), Inches(img_y),
                    Inches(img_w), Inches(img_h),
                )
            except Exception as e:
                print(f"WARNING: Failed to insert image '{image_path}': {e}",
                      file=sys.stderr)
                self._add_image_placeholder(slide, img_x, img_y, img_w, img_h)
        else:
            if image_path:
                print(f"WARNING: Image not found: {image_path}", file=sys.stderr)
            self._add_image_placeholder(slide, img_x, img_y, img_w, img_h)

        # Text side
        text_x = 0.75 if image_side == "right" else (img_x + img_w + 0.5)
        text_w = (13.333 - text_x - 0.75)

        # Bullets
        bullets = conf.get("bullets", [])
        if bullets:
            tb2 = self._add_text_box(
                slide, None,
                Inches(text_x), Inches(img_y),
                Inches(text_w), Inches(img_h),
            )
            first = True
            for b in bullets:
                if first:
                    self._set_text(tb2.text_frame, f"  {b}", "text",
                                   18, False, PP_ALIGN.LEFT, spacing_pt=14)
                    first = False
                else:
                    self._add_para(tb2.text_frame, f"  {b}", "text",
                                   18, False, spacing_pt=14)

        self._add_footer(slide)

    def _add_image_placeholder(self, slide, x, y, w, h):
        """Placeholder rectangle for missing images."""
        ph = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        self._shape_fill(ph, "panel")
        self._shape_line(ph, "border", 1)
        plabel = self._add_text_box(
            slide, None,
            Inches(x + w / 2 - 1), Inches(y + h / 2 - 0.2),
            Inches(2), Inches(0.4),
        )
        self._set_text(plabel.text_frame, "[ Image ]", "muted",
                       13, False, PP_ALIGN.CENTER)

    def _slide_comparison(self, conf):
        """Two-column comparison (before/after)."""
        slide = self._blank_slide()

        # Title
        if conf.get("title"):
            tb = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.6),
            )
            self._set_text(tb.text_frame, conf["title"], "text", 30, True)

        # Two panels
        labels = conf.get("labels", ["Before", "After"])
        left_data = conf.get("left", {})
        right_data = conf.get("right", {})

        col_w = 5.5
        col_h = 4.5
        col_y = 1.3

        for side, data, label in [
            ("left", left_data, labels[0] if len(labels) > 0 else "Before"),
            ("right", right_data, labels[1] if len(labels) > 1 else "After"),
        ]:
            x = 0.75 if side == "left" else (13.333 - 0.75 - col_w)

            # Panel bg
            panel = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(col_y),
                Inches(col_w), Inches(col_h),
            )
            is_primary_side = (side == "right")
            self._shape_fill(panel, "primary" if is_primary_side else "panel")
            self._shape_line(panel, None)

            # Label tag
            label_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.3), Inches(col_y + 0.3),
                Inches(1.2), Inches(0.35),
            )
            self._shape_fill(label_bg, "accent" if is_primary_side else "primary")
            self._shape_line(label_bg, None)
            txt_clr = "text" if is_primary_side else "text_on_primary"
            self._set_text(label_bg.text_frame, label.upper(),
                           txt_clr, 11, True, PP_ALIGN.CENTER)
            label_bg.text_frame.margin_top = Inches(0.05)

            # Content title
            ct = data.get("title", "")
            fill_clr = "primary" if is_primary_side else "panel"
            if ct:
                tb2 = self._add_text_box(
                    slide, fill_clr,
                    Inches(x + 0.3), Inches(col_y + 0.9),
                    Inches(col_w - 0.6), Inches(0.4),
                )
                tc = "text_on_primary" if is_primary_side else "text"
                self._set_text(tb2.text_frame, ct, tc, 20, True)

            # Content body
            cb = data.get("body", "")
            if cb:
                tb3 = self._add_text_box(
                    slide, fill_clr,
                    Inches(x + 0.3), Inches(col_y + 1.4),
                    Inches(col_w - 0.6), Inches(col_h - 1.7),
                )
                tc = "text_on_primary" if is_primary_side else "muted"
                self._set_text(tb3.text_frame, cb, tc, 15, False)

        self._add_footer(slide)

    def _slide_process(self, conf):
        """Horizontal timeline with numbered steps."""
        slide = self._blank_slide()
        steps = conf.get("steps", [])
        n = len(steps)
        if n < 2:
            return self._slide_content(conf)

        # Title
        if conf.get("title"):
            tb = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.6),
            )
            self._set_text(tb.text_frame, conf["title"], "text", 30, True)

        # Timeline layout
        step_w = 2.4
        gap = 0.5
        total = n * step_w + (n - 1) * gap
        start_x = (13.333 - total) / 2
        line_y = 3.85

        # Connecting line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x), Inches(line_y),
            Inches(total), Inches(0.03),
        )
        self._shape_fill(line, "border")
        self._shape_line(line, None)

        for i, step_text in enumerate(steps):
            cx = start_x + i * (step_w + gap)

            # Step circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx + step_w / 2 - 0.3), Inches(line_y - 0.25),
                Inches(0.6), Inches(0.6),
            )
            self._shape_fill(circle, "primary")
            self._shape_line(circle, None)
            num_str = f"{i + 1}"
            self._set_text(circle.text_frame, num_str, "text_on_primary",
                           18, True, PP_ALIGN.CENTER)
            circle.text_frame.margin_top = Inches(0.12)

            # Step description
            tb2 = self._add_text_box(
                slide, None,
                Inches(cx), Inches(line_y + 0.6),
                Inches(step_w), Inches(0.6),
            )
            self._set_text(tb2.text_frame, step_text, "text",
                           16, False, PP_ALIGN.CENTER)

        # Note
        note = conf.get("note", "")
        if note:
            tb3 = self._add_text_box(
                slide, None,
                Inches(2), Inches(line_y + 1.5),
                Inches(9.333), Inches(0.4),
            )
            self._set_text(tb3.text_frame, note, "muted", 11, False,
                           PP_ALIGN.CENTER)

        self._add_footer(slide)

    def _slide_table(self, conf):
        """Structured table with styled headers."""
        slide = self._blank_slide()

        # Title
        if conf.get("title"):
            tb = self._add_text_box(
                slide, None,
                Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.6),
            )
            self._set_text(tb.text_frame, conf["title"], "text", 30, True)

        headers = conf.get("headers", [])
        rows = conf.get("rows", [])
        if not headers or not rows:
            return

        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)
        tbl_w = Inches(11.3)
        tbl_h = min(Inches(0.4 * n_rows + 0.3), Inches(5.5))
        tbl_x = Inches(1.0)
        tbl_y = Inches(1.4)

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h,
        )
        table = table_shape.table

        # Set column widths evenly
        col_w = int(Inches(11.3) / n_cols)
        for ci in range(n_cols):
            table.columns[ci].width = col_w

        # Header row
        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = header
            self._shape_fill(cell, "primary")
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = self._rgb("text_on_primary")
                p.font.name = self.palette.get("font_name", "Microsoft YaHei")
                p.alignment = PP_ALIGN.CENTER
            cell.text_frame.margin_top = Inches(0.06)
            cell.text_frame.margin_bottom = Inches(0.06)

        # Data rows
        for ri, row in enumerate(rows):
            bg = "panel" if ri % 2 == 0 else "bg"
            for ci, cell_text in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = str(cell_text)
                self._shape_fill(cell, bg)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = self._rgb("text")
                    p.font.name = self.palette.get("font_name", "Microsoft YaHei")
                    p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                cell.text_frame.margin_top = Inches(0.04)
                cell.text_frame.margin_bottom = Inches(0.04)
                cell.text_frame.margin_left = Inches(0.15)

        self._add_footer(slide)


# ---- CLI -------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Rich-layout PPTX generator"
    )
    parser.add_argument("--config", help="JSON config file path")
    parser.add_argument("--stdin", action="store_true",
                        help="Read JSON config from stdin")
    parser.add_argument("--json", action="store_true",
                        help="Output JSON report")
    args = parser.parse_args()

    if args.stdin:
        config = json.load(sys.stdin)
    elif args.config:
        with open(args.config, "r", encoding="utf-8") as f:
            config = json.load(f)
    else:
        parser.print_help()
        sys.exit(1)

    engine = LayoutEngine(config)
    path = engine.build()
    report = engine.report()

    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"Generated: {path} ({report['slides']} slides, "
              f"palette: {report['palette']})")


if __name__ == "__main__":
    main()
